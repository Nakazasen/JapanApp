"""Book reader service for parsing various book formats."""
import json
from pathlib import Path
from typing import Dict, List, Optional
import ebooklib
try:
    from ebooklib import epub
except ImportError:
    import ebooklib.epub as epub
from bs4 import BeautifulSoup
from docx import Document
import pandas as pd
from pptx import Presentation


class BaseReader:
    """Base class for all file readers."""
    
    def __init__(self, file_path: Path):
        self.file_path = Path(file_path)
        if not self.file_path.exists():
            raise FileNotFoundError(f"File không tồn tại: {file_path}")
    
    def read(self) -> Dict:
        """Read file content and return dictionary."""
        raise NotImplementedError
    
    def get_chapters(self) -> List[Dict]:
        """Get list of chapters."""
        raise NotImplementedError
    
    def extract_text(self, chapter_index: int = 0) -> str:
        """Extract text from specific chapter."""
        raise NotImplementedError
    
    def search(self, query: str) -> List[Dict]:
        """Search in content."""
        results = []
        chapters = self.get_chapters()
        
        query_lower = query.lower()
        for idx, chapter in enumerate(chapters):
            content = chapter.get("content", "").lower()
            if query_lower in content:
                start_pos = content.find(query_lower)
                context_start = max(0, start_pos - 50)
                context_end = min(len(content), start_pos + len(query) + 50)
                context = chapter["content"][context_start:context_end]
                
                results.append({
                    "chapter_index": idx,
                    "position": start_pos,
                    "context": context,
                    "chapter_title": chapter.get("title", f"Chương {idx + 1}")
                })
        
        return results


class EPUBReader(BaseReader):
    """EPUB file reader."""
    
    def __init__(self, file_path: Path):
        super().__init__(file_path)
        self.book = epub.read_epub(str(self.file_path))
        self._chapters = None
    
    def read(self) -> Dict:
        """Read entire EPUB content."""
        title = self.book.get_metadata('DC', 'title')
        author = self.book.get_metadata('DC', 'creator')
        
        chapters = self.get_chapters()
        
        return {
            "title": title[0][0] if title else "Không có tiêu đề",
            "author": author[0][0] if author else "Không rõ tác giả",
            "chapters": chapters,
            "metadata": {
                "format": "EPUB",
                "language": self.book.get_metadata('DC', 'language'),
                "publisher": self.book.get_metadata('DC', 'publisher'),
            }
        }
    
    def get_chapters(self) -> List[Dict]:
        """Get list of chapters."""
        if self._chapters is not None:
            return self._chapters
        
        chapters = []
        for item in self.book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                content = item.get_content().decode('utf-8')
                soup = BeautifulSoup(content, 'html.parser')
                
                title_tag = soup.find(['h1', 'h2', 'h3', 'title'])
                title = title_tag.get_text().strip() if title_tag else "Chương không có tiêu đề"
                
                for script in soup(["script", "style"]):
                    script.decompose()
                
                text = soup.get_text()
                lines = [line.strip() for line in text.splitlines()]
                text = '\n'.join(line for line in lines if line)
                
                images = []
                for img in soup.find_all('img'):
                    src = img.get('src', '')
                    if src:
                        images.append(src)
                
                chapters.append({
                    "title": title,
                    "content": text,
                    "images": images,
                    "html_content": content
                })
        
        self._chapters = chapters
        return chapters
    
    def extract_text(self, chapter_index: int = 0) -> str:
        """Extract text from specific chapter."""
        chapters = self.get_chapters()
        if 0 <= chapter_index < len(chapters):
            return chapters[chapter_index]["content"]
        return ""


class DOCXReader(BaseReader):
    """Word DOCX file reader."""
    
    def __init__(self, file_path: Path):
        super().__init__(file_path)
        self.doc = Document(str(self.file_path))
        self._chapters = None
    
    def read(self) -> Dict:
        """Read entire DOCX content."""
        title = self.doc.core_properties.title or "Không có tiêu đề"
        author = self.doc.core_properties.author or "Không rõ tác giả"
        
        chapters = self.get_chapters()
        
        return {
            "title": title,
            "author": author,
            "chapters": chapters,
            "metadata": {
                "format": "DOCX",
                "created": str(self.doc.core_properties.created),
                "modified": str(self.doc.core_properties.modified),
            }
        }
    
    def get_chapters(self) -> List[Dict]:
        """Get list of chapters (divided by headings)."""
        if self._chapters is not None:
            return self._chapters
        
        chapters = []
        current_chapter = {"title": "Bắt đầu", "content": "", "images": []}
        
        for para in self.doc.paragraphs:
            if para.style.name.startswith('Heading'):
                if current_chapter["content"].strip():
                    chapters.append(current_chapter)
                
                current_chapter = {
                    "title": para.text.strip() or f"Chương {len(chapters) + 1}",
                    "content": "",
                    "images": []
                }
            else:
                if para.text.strip():
                    current_chapter["content"] += para.text + "\n"
        
        if current_chapter["content"].strip():
            chapters.append(current_chapter)
        
        if not chapters:
            full_text = "\n".join([para.text for para in self.doc.paragraphs if para.text.strip()])
            chapters.append({
                "title": "Nội dung",
                "content": full_text,
                "images": []
            })
        
        self._chapters = chapters
        return chapters
    
    def extract_text(self, chapter_index: int = 0) -> str:
        """Extract text from specific chapter."""
        chapters = self.get_chapters()
        if 0 <= chapter_index < len(chapters):
            return chapters[chapter_index]["content"]
        return ""


class XLSXReader(BaseReader):
    """Excel XLSX file reader."""
    
    def __init__(self, file_path: Path):
        super().__init__(file_path)
        self._chapters = None
    
    def read(self) -> Dict:
        """Read entire XLSX content."""
        excel_file = pd.ExcelFile(str(self.file_path))
        sheet_names = excel_file.sheet_names
        
        chapters = self.get_chapters()
        
        return {
            "title": self.file_path.stem,
            "author": "Không rõ tác giả",
            "chapters": chapters,
            "metadata": {
                "format": "XLSX",
                "sheets": sheet_names,
            }
        }
    
    def get_chapters(self) -> List[Dict]:
        """Get list of sheets as chapters."""
        if self._chapters is not None:
            return self._chapters
        
        chapters = []
        excel_file = pd.ExcelFile(str(self.file_path))
        
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(excel_file, sheet_name=sheet_name)
            
            content_parts = []
            
            if not df.empty:
                content_parts.append(" | ".join(str(col) for col in df.columns))
                content_parts.append("-" * 50)
                
                for idx, row in df.iterrows():
                    row_text = " | ".join(str(val) for val in row.values)
                    content_parts.append(row_text)
            
            content = "\n".join(content_parts)
            
            chapters.append({
                "title": f"Sheet: {sheet_name}",
                "content": content,
                "images": []
            })
        
        self._chapters = chapters
        return chapters
    
    def extract_text(self, chapter_index: int = 0) -> str:
        """Extract text from specific sheet."""
        chapters = self.get_chapters()
        if 0 <= chapter_index < len(chapters):
            return chapters[chapter_index]["content"]
        return ""


class PPTXReader(BaseReader):
    """PowerPoint PPTX file reader."""
    
    def __init__(self, file_path: Path):
        super().__init__(file_path)
        self.prs = Presentation(str(self.file_path))
        self._chapters = None
    
    def read(self) -> Dict:
        """Read entire PPTX content."""
        title = self.file_path.stem
        author = self.prs.core_properties.author if hasattr(self.prs, 'core_properties') else "Không rõ tác giả"
        
        chapters = self.get_chapters()
        
        return {
            "title": title,
            "author": author,
            "chapters": chapters,
            "metadata": {
                "format": "PPTX",
                "slide_count": len(self.prs.slides),
            }
        }
    
    def get_chapters(self) -> List[Dict]:
        """Get list of slides as chapters."""
        if self._chapters is not None:
            return self._chapters
        
        chapters = []
        
        for idx, slide in enumerate(self.prs.slides):
            content_parts = []
            
            title = ""
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.has_text_frame:
                    if shape.is_placeholder and shape.placeholder_format.idx == 0:
                        title = shape.text
                    else:
                        content_parts.append(shape.text)
            
            if not title:
                title = f"Slide {idx + 1}"
            
            content = "\n".join(content_parts)
            
            chapters.append({
                "title": title,
                "content": content,
                "images": []
            })
        
        self._chapters = chapters
        return chapters
    
    def extract_text(self, chapter_index: int = 0) -> str:
        """Extract text from specific slide."""
        chapters = self.get_chapters()
        if 0 <= chapter_index < len(chapters):
            return chapters[chapter_index]["content"]
        return ""


class MOBIReader(BaseReader):
    """Mobipocket/MOBI file reader."""
    
    def __init__(self, file_path: Path):
        super().__init__(file_path)
        self._chapters = None
        self.temp_dir = None
        self.html_path = None
        self.html_files = []  # List of all HTML/XHTML files to read
        # Extract MOBI to HTML and read directly
        try:
            import mobi
            import os
            
            print(f"[MOBIReader] Attempting to extract MOBI file: {self.file_path}")
            
            # Extract MOBI file - returns (temp_dir, html_path) tuple
            extract_result = mobi.extract(str(self.file_path))
            print(f"[MOBIReader] Extract result: {extract_result}")
            
            if extract_result and isinstance(extract_result, tuple) and len(extract_result) == 2:
                self.temp_dir, self.html_path = extract_result
                print(f"[MOBIReader] Temp dir: {self.temp_dir}")
                print(f"[MOBIReader] HTML path: {self.html_path}")
                
                # Check if HTML file exists
                if self.html_path and os.path.exists(self.html_path):
                    print(f"[MOBIReader] HTML file exists, ready to read")
                else:
                    print(f"[MOBIReader] Warning: HTML file does not exist: {self.html_path}")
                    self.html_path = None
            else:
                print(f"[MOBIReader] Warning: mobi.extract returned unexpected result: {extract_result}")
                self.html_path = None
        except ImportError as e:
            print(f"[MOBIReader] Error: mobi library not installed: {e}")
            print("[MOBIReader] Install with: pip install mobi")
            self.html_path = None
        except Exception as e:
            import traceback
            print(f"[MOBIReader] Error: Could not read MOBI file: {e}")
            print(f"[MOBIReader] Traceback: {traceback.format_exc()}")
            self.html_path = None
    
    def read(self) -> Dict:
        """Read entire MOBI content."""
        if self.html_path:
            try:
                import os
                
                print(f"[MOBIReader] Reading HTML file: {self.html_path}")
                
                # Check if file exists
                if not os.path.exists(self.html_path):
                    print(f"[MOBIReader] HTML file does not exist: {self.html_path}")
                    raise FileNotFoundError(f"HTML file not found: {self.html_path}")
                
                # Read file as binary first to check its type
                with open(self.html_path, 'rb') as f:
                    raw_data = f.read()
                
                # Check if file is binary (ZIP signature) - might be wrong file
                if raw_data.startswith(b'PK'):
                    print(f"[MOBIReader] File appears to be a ZIP archive, searching for HTML/XHTML files in temp directory...")
                    # Try to find HTML/XHTML files in the temp directory
                    # MOBI extracted to EPUB structure usually has files in OEBPS/Text/ directory
                    html_files = []
                    if self.temp_dir and os.path.exists(self.temp_dir):
                        # Look for OEBPS/Text/ directory first (EPUB structure)
                        oebps_text_dir = os.path.join(self.temp_dir, 'OEBPS', 'Text')
                        if os.path.exists(oebps_text_dir):
                            print(f"[MOBIReader] Found OEBPS/Text directory: {oebps_text_dir}")
                            search_dirs = [oebps_text_dir]
                        else:
                            # Search entire temp directory
                            search_dirs = [self.temp_dir]
                        
                        for search_dir in search_dirs:
                            for root, dirs, files in os.walk(search_dir):
                                for file in files:
                                    if file.endswith(('.html', '.htm', '.xhtml')):
                                        full_path = os.path.join(root, file)
                                        # Skip cover pages and TOC files
                                        if 'cover' in file.lower() or 'toc' in file.lower():
                                            continue
                                        # Check if it's actually HTML with content
                                        try:
                                            with open(full_path, 'rb') as test_f:
                                                test_data = test_f.read(2048)
                                                if b'<' in test_data and b'>' in test_data:
                                                    # Check if it has actual text content (not just structure)
                                                    decoded_test = test_data.decode('utf-8', errors='ignore')
                                                    if len(decoded_test.replace(' ', '').replace('\n', '')) > 100:
                                                        html_files.append(full_path)
                                        except:
                                            pass
                    
                    if html_files:
                        # Sort files to read in order (part0000.xhtml, part0001.xhtml, etc.)
                        html_files.sort()
                        print(f"[MOBIReader] Found {len(html_files)} HTML/XHTML file(s)")
                        # Store all HTML files for reading
                        self.html_files = html_files
                        # Use the first file for initial reading
                        self.html_path = html_files[0]
                        with open(self.html_path, 'rb') as f:
                            raw_data = f.read()
                    else:
                        raise ValueError("No valid HTML/XHTML files found in extracted MOBI")
                else:
                    # Single HTML file
                    self.html_files = [self.html_path]
                
                # Try multiple encodings to decode the HTML
                html_content = None
                encodings_to_try = ['utf-8', 'latin-1', 'cp1252', 'gbk', 'big5', 'utf-16', 'utf-16-le', 'utf-16-be']
                
                for enc in encodings_to_try:
                    try:
                        decoded = raw_data.decode(enc, errors='strict')
                        # Check if it looks like HTML
                        if '<html' in decoded.lower() or '<body' in decoded.lower() or '<div' in decoded.lower() or '<p' in decoded.lower():
                            html_content = decoded
                            print(f"[MOBIReader] Successfully decoded with encoding: {enc}")
                            break
                    except (UnicodeDecodeError, LookupError):
                        continue
                
                if html_content is None:
                    # Last resort: use errors='replace' with UTF-8
                    html_content = raw_data.decode('utf-8', errors='replace')
                    print(f"[MOBIReader] Using UTF-8 with error replacement (content may be corrupted)")
                
                # Check if content looks like HTML
                if not ('<' in html_content and '>' in html_content):
                    print(f"[MOBIReader] Warning: Content does not look like HTML, attempting to extract text anyway")
                
                # Read all HTML files if multiple files exist
                all_html_content = []
                if hasattr(self, 'html_files') and len(self.html_files) > 1:
                    print(f"[MOBIReader] Reading {len(self.html_files)} HTML/XHTML files...")
                    for html_file in self.html_files:
                        try:
                            with open(html_file, 'rb') as f:
                                file_data = f.read()
                            # Try to decode
                            for enc in ['utf-8', 'latin-1', 'cp1252']:
                                try:
                                    file_content = file_data.decode(enc, errors='replace')
                                    if '<' in file_content and '>' in file_content:
                                        all_html_content.append(file_content)
                                        break
                                except:
                                    continue
                        except Exception as e:
                            print(f"[MOBIReader] Warning: Could not read {html_file}: {e}")
                            continue
                else:
                    # Single file
                    all_html_content = [html_content]
                
                # Parse all HTML content
                all_text_parts = []
                title = self.file_path.stem
                author = "Không rõ tác giả"
                
                for html_content_part in all_html_content:
                    soup = BeautifulSoup(html_content_part, 'html.parser')
                    
                    # Extract title from first file
                    if not title or title == self.file_path.stem:
                        title_tag = soup.find('title')
                        if title_tag:
                            title = title_tag.get_text().strip()
                    
                    # Extract author from first file
                    if author == "Không rõ tác giả":
                        author_tag = soup.find('meta', {'name': 'author'})
                        if author_tag:
                            author = author_tag.get('content', '')
                    
                    # Remove script and style tags
                    for script in soup(["script", "style"]):
                        script.decompose()
                    
                    # Extract text content
                    text_content = soup.get_text()
                    lines = [line.strip() for line in text_content.splitlines()]
                    text_content = '\n'.join(line for line in lines if line)
                    
                    if text_content and len(text_content.strip()) > 10:  # Only add if has meaningful content
                        all_text_parts.append(text_content)
                
                # Combine all text parts
                combined_text = '\n\n'.join(all_text_parts)
                
                # Split into chapters
                chapters = []
                if combined_text:
                    # Try to split by headings if we have multiple HTML files
                    if len(all_html_content) > 1:
                        # Each HTML file becomes a chapter
                        for idx, html_content_part in enumerate(all_html_content):
                            soup = BeautifulSoup(html_content_part, 'html.parser')
                            for script in soup(["script", "style"]):
                                script.decompose()
                            
                            # Find title in this file
                            chapter_title = f"Chương {idx + 1}"
                            h1_tag = soup.find('h1')
                            if h1_tag:
                                chapter_title = h1_tag.get_text().strip()
                            else:
                                h2_tag = soup.find('h2')
                                if h2_tag:
                                    chapter_title = h2_tag.get_text().strip()
                            
                            text_content = soup.get_text()
                            lines = [line.strip() for line in text_content.splitlines()]
                            text_content = '\n'.join(line for line in lines if line)
                            
                            if text_content and len(text_content.strip()) > 10:
                                chapters.append({
                                    "title": chapter_title,
                                    "content": text_content,
                                    "html_content": html_content_part,
                                    "images": []
                                })
                    else:
                        # Single file - try to split by headings
                        soup = BeautifulSoup(html_content, 'html.parser')
                        headings = soup.find_all(['h1', 'h2', 'h3'])
                        if headings:
                            current_chapter = []
                            current_title = "Bắt đầu"
                            
                            for element in soup.find_all(['h1', 'h2', 'h3', 'p', 'div']):
                                if element.name in ['h1', 'h2', 'h3']:
                                    if current_chapter:
                                        chapter_text = '\n'.join(current_chapter).strip()
                                        if chapter_text:
                                            chapters.append({
                                                "title": current_title,
                                                "content": chapter_text,
                                                "html_content": "",
                                                "images": []
                                            })
                                    current_title = element.get_text().strip() or f"Chương {len(chapters) + 1}"
                                    current_chapter = []
                                else:
                                    text = element.get_text().strip()
                                    if text:
                                        current_chapter.append(text)
                            
                            if current_chapter:
                                chapter_text = '\n'.join(current_chapter).strip()
                                if chapter_text:
                                    chapters.append({
                                        "title": current_title,
                                        "content": chapter_text,
                                        "html_content": "",
                                        "images": []
                                    })
                        
                        # If no headings found, create single chapter
                        if not chapters:
                            chapters.append({
                                "title": "Nội dung",
                                "content": combined_text if combined_text else text_content,
                                "html_content": html_content,
                                "images": []
                            })
                
                print(f"[MOBIReader] Successfully read MOBI, chapters: {len(chapters)}")
                
                return {
                    "title": title,
                    "author": author,
                    "chapters": chapters,
                    "metadata": {"format": "MOBI"}
                }
            except Exception as e:
                import traceback
                print(f"[MOBIReader] Error reading HTML from MOBI: {e}")
                print(f"[MOBIReader] Traceback: {traceback.format_exc()}")
        
        # Fallback: return basic info with warning
        print(f"[MOBIReader] Falling back to error message")
        return {
            "title": self.file_path.stem,
            "author": "Không rõ tác giả",
            "chapters": [{
                "title": "Nội dung",
                "content": "⚠️ Không thể đọc file MOBI. Vui lòng chuyển đổi sang EPUB hoặc DOCX để đọc tốt hơn.\n\nĐể chuyển đổi MOBI sang EPUB, bạn có thể sử dụng:\n- Calibre (miễn phí, khuyến nghị)\n- Online converters\n- Amazon Kindle tools",
                "images": []
            }],
            "metadata": {"format": "MOBI"}
        }
    
    def get_chapters(self) -> List[Dict]:
        """Get list of chapters."""
        if self._chapters is not None:
            return self._chapters
        
        # Read book data to get chapters
        book_data = self.read()
        self._chapters = book_data.get("chapters", [])
        return self._chapters
    
    def extract_text(self, chapter_index: int = 0) -> str:
        """Extract text from specific chapter."""
        chapters = self.get_chapters()
        if 0 <= chapter_index < len(chapters):
            return chapters[chapter_index]["content"]
        return ""
    
    def __del__(self):
        """Cleanup extracted temp files if exists."""
        try:
            import os
            import shutil
            if self.temp_dir and os.path.exists(self.temp_dir):
                # Note: mobi.extract creates temp files that should be cleaned up
                # but we'll leave them for now to avoid issues with file locks
                # In production, you might want to implement a cleanup mechanism
                pass
        except:
            pass


class AZW3Reader(MOBIReader):
    """AZW3 (Kindle Format 8) file reader.
    
    AZW3 is similar to MOBI but uses EPUB-like structure.
    We can use the same extraction method as MOBI.
    """
    
    def __init__(self, file_path: Path):
        # AZW3 uses the same extraction method as MOBI
        super().__init__(file_path)
    
    def read(self) -> Dict:
        """Read entire AZW3 content."""
        # Use parent MOBIReader's read method but update metadata
        result = super().read()
        if result.get("metadata"):
            result["metadata"]["format"] = "AZW3"
        return result


class BookReaderService:
    """Service for reading books in various formats."""
    
    SUPPORTED_FORMATS = {
        ".epub": EPUBReader,
        ".mobi": MOBIReader,
        ".prc": MOBIReader,  # PRC is also Mobipocket format
        ".azw3": AZW3Reader,  # AZW3 is Kindle Format 8
        ".docx": DOCXReader,
        ".doc": DOCXReader,
        ".xlsx": XLSXReader,
        ".xls": XLSXReader,
        ".pptx": PPTXReader,
        ".ppt": PPTXReader,
    }
    
    @staticmethod
    def create_reader(file_path: Path) -> BaseReader:
        """Create appropriate reader for file format."""
        file_path = Path(file_path)
        suffix = file_path.suffix.lower()
        
        if suffix not in BookReaderService.SUPPORTED_FORMATS:
            raise ValueError(f"Định dạng file không được hỗ trợ: {suffix}")
        
        reader_class = BookReaderService.SUPPORTED_FORMATS[suffix]
        return reader_class(file_path)
    
    @staticmethod
    def read_book(file_path: Path) -> Dict:
        """Read book and return parsed data."""
        reader = BookReaderService.create_reader(file_path)
        return reader.read()
    
    @staticmethod
    def get_chapters(file_path: Path) -> List[Dict]:
        """Get chapters from book."""
        reader = BookReaderService.create_reader(file_path)
        return reader.get_chapters()
    
    @staticmethod
    def search_in_book(file_path: Path, query: str) -> List[Dict]:
        """Search in book content."""
        reader = BookReaderService.create_reader(file_path)
        return reader.search(query)

